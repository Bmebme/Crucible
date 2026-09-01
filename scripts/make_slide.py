# -*- coding: utf-8 -*-
"""生成《AI 辅助漏洞验证的知识库》PPT（2 页）。
第 1 页: 知识库的特点与作用（正面介绍，不含技术对比）
第 2 页: 知识库架构（双引擎供给 · MCP 出口 · 验证回写）
配色沿用 dataviz 验证色板（浅色模式）。重新生成: python3 scripts/make_slide.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

INK    = RGBColor(0x0B, 0x0B, 0x0B)
INK2   = RGBColor(0x52, 0x51, 0x4E)
MUTED  = RGBColor(0x89, 0x87, 0x81)
HAIR   = RGBColor(0xE1, 0xE0, 0xD9)
SURF   = RGBColor(0xFC, 0xFC, 0xFB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
AQUA   = RGBColor(0x1B, 0xAF, 0x7A)
YELLOW = RGBColor(0xED, 0xA1, 0x00)
VIOLET = RGBColor(0x4A, 0x3A, 0xA7)
GREEN  = RGBColor(0x00, 0x83, 0x00)
GOOD   = RGBColor(0x0C, 0xA3, 0x0C)
SERIOUS= RGBColor(0xEC, 0x83, 0x5A)
GREEN_BG = RGBColor(0xE8, 0xF3, 0xEA)
TINT = {
    "blue":   RGBColor(0xE8, 0xF0, 0xFB),
    "aqua":   RGBColor(0xE6, 0xF5, 0xEE),
    "violet": RGBColor(0xEC, 0xEA, 0xF6),
}

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    try:
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = SURF
    except Exception:
        pass
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, round_=False, radius=0.08):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def arrow(s, x, y, w, h, fill, shape=MSO_SHAPE.RIGHT_ARROW):
    shp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def set_run(run, text, size, bold=False, color=INK):
    run.text = text
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = FONT
    rPr = run._r.get_or_add_rPr()
    if rPr.find(qn('a:ea')) is None:
        ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT})
        rPr.append(ea)


def text(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=2, line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, runs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, size, bold, color) in runs:
            set_run(p.add_run(), t, size, bold, color)
    return tb


# ============================================================ 页 1：特点与作用
s1 = new_slide()
text(s1, 0.5, 0.14, 12.33, 0.5, [[("AI 辅助漏洞验证的知识库：特点与作用", 26, True, INK)]])
# 与传统问答知识库的对比带
box(s1, 0.5, 0.72, 6.1, 0.56, fill=HAIR, round_=True, radius=0.12)
text(s1, 0.72, 0.77, 5.7, 0.46,
     [[("✗ 传统问答知识库：回答「产品包含哪些模块」—— 结构盘点", 10.5, True, INK2)],
      [("段落 · 静态 · 问题驱动", 9, False, MUTED)]], space_after=1)
box(s1, 6.73, 0.72, 6.1, 0.56, fill=GREEN_BG, line=GREEN, line_w=1.0, round_=True, radius=0.12)
text(s1, 6.95, 0.77, 5.7, 0.46,
     [[("✓ 本方案：回答「模块的通信与安全机制、面对哪些威胁」", 10.5, True, INK)],
      [("字段化知识点 · 只供知识，推断归 Agent · 验证回写", 9, False, INK2)]], space_after=1)

CARDS = [
    (0.5, BLUE, "blue", "① 场景构建 · 发散",
     ["场景点：用户行为的完整路径——入口（从哪进）、组件链（经过哪些模块）、数据流（落到哪）、信任变化（信任怎么变）",
      "威胁点：攻击目标的可选项——目标层级、所需能力、失败降级备选",
      "经验模式：利用条件——前提（版本/鉴权/可达性）+ 探测方法 + 溯源",
      "案例库：同类产品案例，按当前版本过滤"],
     ["无需二次提炼：直接消费前提与探测字段，少一轮失真",
      "提前过滤：前提不符（版本/不可达）不进候选",
      "直接可用：带步骤与预期，直接进路径规划"],
     "AttackScenario 候选集（Agent 筛选与排序）"),
    (4.65, AQUA, "aqua", "② 路径规划 · 收敛",
     ["机制点：模块的实现真相——通信机制（协议/调用方式/数据流向）+ 安全机制（鉴权/校验/防护），回答「怎么走、哪里有防护」",
      "历史问题：路径相关模块出过的问题与验证记录——哪些坑被踩过、结果如何",
      "关联召回：知识点间 wikilink 关联，多跳路径的中间机制点不漏召回"],
     ["机制点不漏：Agent 拿全所需的点再组织路径",
      "历史提示：据历史问题评估路径可行性",
      "证据后置：由第③步实测验证得出"],
     "机制点集 + 来源标注（路径推断由 Agent 完成）"),
    (8.8, VIOLET, "violet", "③ POC 生成与验证 · 执行",
     ["历史 POC：验证过的可复用骨架——占位符 + 请求模板 + 验证断言，实例化即用",
      "历史结果：本产品实测记录——哪个 payload 被拦、哪个绕过成功、哪个是误报",
      "状态字段：每条知识带验证状态——成功加权复用 / 拦截成负知识 / 误报修正前提"],
     ["不从零试：站在历史 POC 上实例化，不重复撞墙",
      "复利变准：验证回写 → 状态流转 → 检索重排序"],
     "历史 POC 模板 + 验证记录（执行与判定归 Agent）"),
]
CARD_W, CARD_Y, CARD_H = 4.03, 1.42, 3.9
for cx, color, tint_key, title, items, effects, out_line in CARDS:
    box(s1, cx, CARD_Y, CARD_W, CARD_H, fill=WHITE, line=color, line_w=1.5, round_=True, radius=0.04)
    box(s1, cx, CARD_Y, CARD_W, 0.4, fill=color)
    text(s1, cx + 0.14, CARD_Y + 0.03, CARD_W - 0.28, 0.34,
         [[(title, 13, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s1, cx + 0.16, CARD_Y + 0.48, CARD_W - 0.32, 0.22,
         [[("提供 —— 关键知识点", 10.5, True, color)]])
    text(s1, cx + 0.16, CARD_Y + 0.72, CARD_W - 0.32, 1.25,
         [[("• " + t, 10, False, INK2)] for t in items], space_after=2)
    box(s1, cx + 0.12, CARD_Y + 1.82, CARD_W - 0.24, 1.2, fill=TINT[tint_key], round_=True, radius=0.06)
    text(s1, cx + 0.26, CARD_Y + 1.9, CARD_W - 0.52, 1.05,
         [[("作用", 9.5, True, color)]] + [[("• " + e, 9.5, False, INK2)] for e in effects],
         space_after=2, line_spacing=1.05)
    text(s1, cx + 0.16, CARD_Y + 3.14, CARD_W - 0.32, 0.4,
         [[("输出：", 10, True, color), (out_line, 10, False, INK2)]])

# 知识点字段示例带
box(s1, 0.5, 5.44, 12.33, 0.4, fill=WHITE, line=HAIR, line_w=1.0, round_=True, radius=0.5)
text(s1, 0.5, 5.44, 12.33, 0.4,
     [[("知识点长这样「经验模式 · ImageTragick」：", 9.5, True, VIOLET),
       ("前提 = convert < 6.9.3-9 ｜ 探测 = mvg payload ｜ 状态 = verified_success（被拦 → msl 绕过）｜ 溯源 = CVE-2016-3714",
        9.5, False, INK2)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 特点（八条，带说明）
text(s1, 0.5, 5.94, 4.0, 0.22, [[("知识库的特点（八条）", 11, True, MUTED)]])
FEATURES = [
    ("可执行性", "知识带步骤与工具参数，Agent 直接消费", BLUE),
    ("条件化", "每条带前提（版本/鉴权/可达），不符不召回", ORANGE),
    ("验证状态机", "成功加权 · 拦截负知识 · 误报修正前提", AQUA),
    ("失败知识入库", "拦截特征与误报记录，防重复撞墙", YELLOW),
    ("经验复利", "验证回写，知识随实测越用越准", GREEN),
    ("目标坐标系", "目标层级 + 降级链，支持动态调整", VIOLET),
    ("召回保证", "关联结构使多跳路径中间点不漏", GOOD),
    ("版本敏感", "知识按产品版本分叉与过滤", SERIOUS),
]
for i, (label, desc, color) in enumerate(FEATURES):
    row, col = divmod(i, 4)
    cx = 0.5 + col * 3.11
    cy = 6.18 + row * 0.6
    box(s1, cx, cy, 3.0, 0.55, fill=WHITE, line=HAIR, line_w=1.0, round_=True, radius=0.08)
    box(s1, cx, cy, 0.07, 0.55, fill=color)
    text(s1, cx + 0.14, cy + 0.03, 2.8, 0.22, [[(label, 10, True, INK)]])
    text(s1, cx + 0.14, cy + 0.27, 2.8, 0.26, [[(desc, 8.5, False, MUTED)]])

# ============================================================ 页 2：架构
s2 = new_slide()
text(s2, 0.5, 0.16, 12.33, 0.55, [[("知识库架构：双引擎供给 · MCP 出口 · 验证回写", 26, True, INK)]])

# 知识来源
box(s2, 0.5, 1.0, 12.33, 0.55, fill=HAIR, round_=True, radius=0.5)
text(s2, 0.5, 1.0, 12.33, 0.55,
     [[("知识来源：产品文档 · API 规范 · 源代码 · 漏洞案例 · 历史验证报告", 12, True, INK2)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 双引擎
box(s2, 0.5, 1.78, 6.0, 1.5, fill=WHITE, line=ORANGE, line_w=1.5, round_=True, radius=0.06)
box(s2, 0.5, 1.78, 6.0, 0.4, fill=ORANGE)
text(s2, 0.64, 1.81, 5.7, 0.34, [[("llm_wiki · 叙述层（知识库主体）", 13, True, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)
text(s2, 0.66, 2.3, 5.7, 0.95,
     [[("• 页面目录：场景 / 威胁 / 经验 / 历史 POC / 验证记录", 10.5, False, INK2)],
      [("• 检索：关键词 + 向量 + 一跳图（混合召回）", 10.5, False, INK2)],
      [("• 演化：验证结果回写 → 状态流转 → 检索加权", 10.5, False, INK2)]], space_after=3)

box(s2, 6.83, 1.78, 6.0, 1.5, fill=WHITE, line=AQUA, line_w=1.5, round_=True, radius=0.06)
box(s2, 6.83, 1.78, 6.0, 0.4, fill=AQUA)
text(s2, 6.97, 1.81, 5.7, 0.34, [[("LightRAG · 实体层（v1 可选 · 召回增强）", 13, True, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)
text(s2, 6.99, 2.3, 5.7, 0.95,
     [[("• 实体类型：端点 / 服务 / 参数 / 目标 / 能力", 10.5, False, INK2)],
      [("• 知识点增多后出现多跳漏召回时按需引入", 10.5, False, INK2)],
      [("• 只做召回增强，不做路径组织", 10.5, False, INK2)]], space_after=3)

arrow(s2, 2.95, 1.57, 0.18, 0.19, MUTED, shape=MSO_SHAPE.DOWN_ARROW)
arrow(s2, 9.78, 1.57, 0.18, 0.19, MUTED, shape=MSO_SHAPE.DOWN_ARROW)

# MCP 出口
box(s2, 0.5, 3.52, 12.33, 0.85, fill=WHITE, line=VIOLET, line_w=1.5, round_=True, radius=0.1)
text(s2, 0.5, 3.58, 12.33, 0.75,
     [[("知识出口 MCP（工具面按步骤切片）", 12.5, True, INK)],
      [("kb_scenario 场景候选 · kb_goal 目标与降级备选 · kb_poc 历史 POC · kb_search 兜底检索 · kb_record_verification 验证回写",
        10, False, INK2)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=1)
arrow(s2, 3.45, 3.3, 0.18, 0.2, VIOLET, shape=MSO_SHAPE.DOWN_ARROW)
arrow(s2, 9.78, 3.3, 0.18, 0.2, VIOLET, shape=MSO_SHAPE.DOWN_ARROW)

# 三步流水线
STEPS2 = [("① 场景构建", "供场景点 · 威胁点 · 经验", 0.5),
          ("② 路径规划", "供机制点 · 历史问题", 4.75),
          ("③ POC 生成 + 实测", "供历史 POC · 验证记录", 9.0)]
for title, sub, sx in STEPS2:
    box(s2, sx, 4.62, 3.83, 0.95, fill=WHITE, line=BLUE, line_w=1.5, round_=True, radius=0.1)
    text(s2, sx + 0.12, 4.7, 3.6, 0.8,
         [[(title, 12.5, True, INK)], [(sub, 9.5, False, MUTED)]], space_after=2)
arrow(s2, 2.36, 4.39, 0.18, 0.2, VIOLET, shape=MSO_SHAPE.DOWN_ARROW)
arrow(s2, 6.62, 4.39, 0.18, 0.2, VIOLET, shape=MSO_SHAPE.DOWN_ARROW)
arrow(s2, 10.88, 4.39, 0.18, 0.2, VIOLET, shape=MSO_SHAPE.DOWN_ARROW)
arrow(s2, 4.36, 5.06, 0.36, 0.22, BLUE)
arrow(s2, 8.61, 5.06, 0.36, 0.22, BLUE)

# 验证回写
box(s2, 0.5, 5.92, 12.33, 0.78, fill=GREEN_BG, line=GREEN, line_w=1.0, round_=True, radius=0.08)
text(s2, 0.8, 6.0, 11.7, 0.62,
     [[("验证回写：Agent 实测执行 → 结果回写 → 状态流转（成功加权 / 拦截负知识 / 误报修正）→ 检索重排序",
        12, True, INK)],
      [("实测执行与证据获得归 Agent · 状态维护与检索加权归 KB", 9.5, False, INK2)]], space_after=1)
arrow(s2, 10.88, 5.59, 0.18, 0.31, BLUE, shape=MSO_SHAPE.DOWN_ARROW)

# 回写反馈
conn = s2.shapes.add_connector(MSO_CONNECTOR.ELBOW, Inches(1.1), Inches(5.92), Inches(1.1), Inches(3.28))
conn.line.color.rgb = GREEN
conn.line.width = Pt(1.5)
text(s2, 0.62, 4.42, 2.0, 0.2, [[("验证结果回写", 9, False, MUTED)]])

text(s2, 6.5, 7.1, 6.33, 0.3,
     [[("设计稿 v1.4 · 单引擎 llm_wiki 起步 + 可选 LightRAG 召回增强", 9, False, MUTED)]],
     align=PP_ALIGN.RIGHT)

# ============================================================ 页 3：为什么需要专门的知识库
s3 = new_slide()
text(s3, 0.5, 0.16, 12.33, 0.55, [[("为什么需要专门的知识库（而非通用问答库）", 26, True, INK)]])

ROWS3 = [
    (BLUE, "①", "回答的问题不同", "结构盘点 vs 机制与威胁",
     ["回答「产品包含哪些模块」",
      "模块清单 · 职责描述 · 接口列表",
      "作者视角的结构描述"],
     ["回答「模块怎么通信、有什么安全机制、面对哪些威胁」",
      "通信机制 + 安全机制 + 威胁（三问定级）",
      "攻击者视角"]),
    (VIOLET, "②", "知识形态不同", "段落 vs 字段化知识点",
     ["返回段落：案例文本 · 模式描述 · POC 文本",
      "Agent 使用前须自行二次提炼前提与探测",
      "二次提炼 = 又一轮 LLM 转化，会错会漏会失真"],
     ["字段化知识点：前提 / 探测 / 验证状态",
      "提炼在入库时一次完成：专家预编译 + 人工审核 + 实测修正",
      "Agent 直接消费字段，跳过二次提炼"]),
    (AQUA, "③", "实验反馈源", "无闭环 vs 验证回写",
     ["知识从新文档演化，无实验反馈",
      "「可能注入」永远是「可能」"],
     ["验证回写闭环：成功加权 · 拦截成负知识 · 误报修正前提",
      "知识随验证次数复利变准"]),
]
for i, (color, num, name, tagline, trad, ours) in enumerate(ROWS3):
    ry = 0.95 + i * 1.8
    # 左侧标签
    box(s3, 0.5, ry, 2.3, 1.62, fill=WHITE, line=color, line_w=1.5, round_=True, radius=0.08)
    text(s3, 0.72, ry + 0.14, 1.9, 1.35,
         [[(num, 20, True, color)], [(name, 13, True, INK)], [(tagline, 9, False, MUTED)]],
         space_after=3)
    # 传统问答知识库
    box(s3, 2.95, ry, 4.9, 1.62, fill=HAIR, round_=True, radius=0.08)
    text(s3, 3.13, ry + 0.12, 4.55, 1.4,
         [[("传统问答知识库", 10.5, True, MUTED)]] +
         [[("• " + t, 9.5, False, INK2)] for t in trad], space_after=3)
    # 本方案攻击知识库
    box(s3, 8.0, ry, 4.83, 1.62, fill=GREEN_BG, line=GREEN, line_w=1.0, round_=True, radius=0.08)
    text(s3, 8.18, ry + 0.12, 4.5, 1.4,
         [[("本方案攻击知识库", 10.5, True, GREEN)]] +
         [[("• " + t, 9.5, False, INK2)] for t in ours], space_after=3)

text(s3, 0.5, 6.55, 12.33, 0.4,
     [[("共同根源：攻击验证的消费动作是「行动」，通用问答库的消费动作是「回答」—— 知识库必须按消费动作设计",
        11, True, INK)]], align=PP_ALIGN.CENTER)

OUT = "../slides/AI辅助漏洞验证-知识库作用与特点.pptx"
prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
